"""PPTX reader backed by python-pptx.

Emits a "Slide N" heading per slide, the slide title, body shapes (with
LaTeX-ized OMML equations), and speaker notes.
"""

from __future__ import annotations

from pathlib import Path
from typing import Any

from ..deps import require
from ..engine.model import Document
from ..math_converters.omml import omath_element_to_latex
from .base import Reader

_DRAWING_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
_A_T = f"{{{_DRAWING_NS}}}t"
_MATH_NS = "http://schemas.openxmlformats.org/officeDocument/2006/math"
_M_OMATH = f"{{{_MATH_NS}}}oMath"
_MC_NS = "http://schemas.openxmlformats.org/markup-compatibility/2006"
_MC_ALTERNATE = f"{{{_MC_NS}}}AlternateContent"
_MC_CHOICE = f"{{{_MC_NS}}}Choice"


class PptxReader(Reader):
    """Read .pptx files into a :class:`Document`."""

    extensions = frozenset({".pptx"})
    name = "pptx"

    def read(self, path: Path, **kwargs: Any) -> Document:
        require("pptx", "reading PPTX files")
        import pptx

        prs = pptx.Presentation(path)
        doc = Document()

        for index, slide in enumerate(prs.slides, start=1):
            doc.add_heading(f"Slide {index}", level=2)

            for shape in slide.shapes:
                if shape.has_text_frame:
                    self._parse_text_frame(shape.text_frame, doc)
                elif shape.has_table:
                    self._parse_table(shape.table, doc)

            notes_slide = getattr(slide, "notes_slide", None)
            if notes_slide and notes_slide.notes_text_frame:
                notes = notes_slide.notes_text_frame.text.strip()
                if notes:
                    doc.add_paragraph(f"Notes: {notes}")

        return doc

    def _parse_text_frame(self, tf: Any, doc: Document) -> None:
        for paragraph in tf.paragraphs:
            parts: list[str] = []
            for child in paragraph._p:
                tag = child.tag
                if tag.endswith("}r"):
                    parts.append("".join(t.text or "" for t in child.findall(_A_T)))
                elif tag == _M_OMATH:
                    parts.append(f" ${omath_element_to_latex(child)}$ ")
                elif tag == _MC_ALTERNATE:
                    choice = child.find(_MC_CHOICE)
                    if choice is not None:
                        om = choice.find(_M_OMATH)
                        if om is not None:
                            parts.append(f" ${omath_element_to_latex(om)}$ ")

            text = "".join(parts).strip()
            if text:
                doc.add_paragraph(text)

    def _parse_table(self, table: Any, doc: Document) -> None:
        rows: list[list[str]] = []
        for row in table.rows:
            rows.append([cell.text.strip() for cell in row.cells])
        if rows and any(any(r) for r in rows):
            doc.add_table(rows[1:], header=rows[0])
