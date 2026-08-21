"""Shared fixtures for the tokens.md test suite."""

from __future__ import annotations

import shutil
import tempfile
import zipfile
from pathlib import Path

import pytest

# --- pytest plugins and output configuration --------------------------------
# Install for better output:
#   pip install pytest-sugar pytest-instafail
# Run with:
#   pytest -v --tb=short -ra        # verbose, short tracebacks, summary table
#   pytest -x -v                    # stop at first failure
#   pytest --instafail -v           # show failures as they happen


def pytest_configure(config: pytest.Config) -> None:
    """Enable detailed coverage table when -v is passed; otherwise let --cov-report= suppress the table."""
    if config.option.verbose > 0:
        config.option.cov_report = {"term-missing:skip-covered": None}


def pytest_terminal_summary(terminalreporter: pytest.TerminalReporter, exitstatus: int, config: pytest.Config) -> None:
    """Print a clean single-line coverage percentage whenever coverage is enabled."""
    if getattr(config.option, "cov_source", None):
        try:
            cov_plugin = config.pluginmanager.get_plugin("_cov")
            if cov_plugin and hasattr(cov_plugin, "cov_controller") and cov_plugin.cov_controller:
                import io

                cov = cov_plugin.cov_controller.cov
                stream = io.StringIO()
                total = cov.report(file=stream)
                terminalreporter.write_sep("=", f"Total Coverage: {total:.0f}%", bold=True, green=(total >= 70), yellow=(50 <= total < 70), red=(total < 50))
        except Exception:
            pass


@pytest.fixture
def tmd_workspace(tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> Path:
    """Point the workspace temp dir at a per-test tmp_path."""
    monkeypatch.setattr(tempfile, "gettempdir", lambda: str(tmp_path))
    return tmp_path


@pytest.fixture
def client(tmd_workspace: Path):
    """FastAPI TestClient with an isolated temp workspace."""
    from fastapi.testclient import TestClient

    from backend.app import create_app

    app = create_app()
    with TestClient(app) as test_client:
        yield test_client


@pytest.fixture
def tmp_output(tmp_path: Path) -> Path:
    out = tmp_path / "output"
    out.mkdir()
    return out


@pytest.fixture
def sample_pdf(tmp_path: Path) -> Path:
    """Generate a minimal single-page PDF containing known text."""
    import pymupdf

    pdf_path = tmp_path / "sample.pdf"
    doc = pymupdf.open()
    page = doc.new_page()
    page.insert_text((72, 72), "Hello from tokens.md")
    doc.save(str(pdf_path))
    doc.close()
    return pdf_path


@pytest.fixture
def sample_md(tmp_path: Path) -> Path:
    path = tmp_path / "notes.md"
    path.write_text("# Notes\n\nSome markdown content here.\n", encoding="utf-8")
    return path


@pytest.fixture
def sample_docx(tmp_path: Path) -> Path:
    """Regenerated with python-docx; keeps the strings the existing test asserts."""
    from docx import Document

    doc = Document()
    doc.add_paragraph("Dear friend,")
    doc.add_paragraph("This is a letter.")
    path = tmp_path / "letter.docx"
    doc.save(str(path))
    return path


@pytest.fixture
def sample_docx_headed(tmp_path: Path) -> Path:
    """A DOCX with a real Heading 1 style, for heading-inference tests."""
    from docx import Document

    doc = Document()
    doc.add_heading("Chapter One", level=1)
    doc.add_paragraph("Hello from python-docx.")
    path = tmp_path / "headed.docx"
    doc.save(str(path))
    return path


@pytest.fixture
def sample_pptx(tmp_path: Path) -> Path:
    """A minimal PPTX with one slide, a title, and body text."""
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Hello slides"
    slide.placeholders[1].text = "Body text here."
    path = tmp_path / "deck.pptx"
    prs.save(str(path))
    return path


@pytest.fixture
def sample_xlsx(tmp_path: Path) -> Path:
    """A minimal XLSX with a header row and one data row."""
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    assert ws is not None
    ws.append(["name", "age"])
    ws.append(["alice", 30])
    path = tmp_path / "data.xlsx"
    wb.save(str(path))
    return path


@pytest.fixture
def sample_odt(tmp_path: Path) -> Path:
    """A minimal ODT with one paragraph."""
    from odf.opendocument import OpenDocumentText
    from odf.teletype import addTextToElement
    from odf.text import P

    doc = OpenDocumentText()
    p = P()
    addTextToElement(p, "Hello from ODF")
    doc.text.addElement(p)
    path = tmp_path / "sample.odt"
    doc.save(str(path))
    return path


@pytest.fixture
def sample_rtf(tmp_path: Path) -> Path:
    path = tmp_path / "sample.rtf"
    path.write_text(r"{\rtf1\ansi Hello {\b bold} RTF.}", encoding="utf-8")
    return path


@pytest.fixture
def sample_eml(tmp_path: Path) -> Path:
    path = tmp_path / "sample.eml"
    path.write_text("From: a@example.com\nTo: b@example.com\nSubject: Hello\n\nBody text.\n", encoding="utf-8")
    return path


@pytest.fixture
def sample_srt(tmp_path: Path) -> Path:
    path = tmp_path / "sample.srt"
    path.write_text("1\n00:00:01,000 --> 00:00:03,000\nHello subtitle\n", encoding="utf-8")
    return path


@pytest.fixture
def sample_tex(tmp_path: Path) -> Path:
    path = tmp_path / "sample.tex"
    path.write_text("\\section{Introduction}\nSome body text here.\n", encoding="utf-8")
    return path


@pytest.fixture
def sample_azw4(tmp_path: Path, sample_pdf: Path) -> Path:
    """AZW4 is a PDF wrapper; a copy of sample_pdf with a .azw4 suffix."""
    path = tmp_path / "sample.azw4"
    shutil.copyfile(sample_pdf, path)
    return path


@pytest.fixture
def sample_csv(tmp_path: Path) -> Path:
    path = tmp_path / "data.csv"
    path.write_text("name,age\nalice,30\nbob,25\n", encoding="utf-8")
    return path


@pytest.fixture
def sample_json(tmp_path: Path) -> Path:
    path = tmp_path / "config.json"
    path.write_text('{"host": "localhost", "port": 8080}', encoding="utf-8")
    return path


@pytest.fixture
def sample_txt(tmp_path: Path) -> Path:
    path = tmp_path / "readme.txt"
    path.write_text("Plain text content.\n", encoding="utf-8")
    return path


@pytest.fixture
def sample_html(tmp_path: Path) -> Path:
    path = tmp_path / "page.html"
    path.write_text("<html><body><article><h1>Title</h1><p>Body text here.</p></article></body></html>", encoding="utf-8")
    return path


# --- Math fixtures ---------------------------------------------------------


def _inject_docx_math(tmp_path: Path, omml: str) -> Path:
    """Build a docx with python-docx, then splice OMML into word/document.xml.

    Rewrites the package preserving every other part ([Content_Types].xml,
    rels, styles) so python-docx can still open the result.
    """
    from docx import Document

    doc = Document()
    doc.add_paragraph("Solve for x:")
    path = tmp_path / "math.docx"
    doc.save(str(path))

    tmp = path.with_suffix(".tmp.docx")
    with zipfile.ZipFile(path) as zin, zipfile.ZipFile(tmp, "w") as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename == "word/document.xml":
                xml = data.decode("utf-8")
                omml_xml = f'<w:p><m:oMathPara xmlns:m="http://schemas.openxmlformats.org/officeDocument/2006/math">{omml}</m:oMathPara></w:p>'
                # Insert OMML paragraph AFTER the first paragraph (after first </w:p>)
                xml = xml.replace("</w:p>", "</w:p>" + omml_xml, 1)
                data = xml.encode("utf-8")
            zout.writestr(item, data)
    shutil.move(tmp, path)
    return path


@pytest.fixture
def sample_docx_math(tmp_path: Path) -> Path:
    # x^2 + 1 as OMML: <m:sSup><m:e><m:r><m:t>x</m:t></m:r></m:e><m:sup><m:r><m:t>2</m:t></m:r></m:sup></m:sSup>
    omml = "<m:sSup><m:e><m:r><m:t>x</m:t></m:r></m:e><m:sup><m:r><m:t>2</m:t></m:r></m:sup></m:sSup><m:r><m:t>+1</m:t></m:r>"
    return _inject_docx_math(tmp_path, omml)


@pytest.fixture
def sample_tex_math(tmp_path: Path) -> Path:
    path = tmp_path / "math.tex"
    path.write_text("\\section{Derivation}\nThe energy is $E = mc^2$.\n\\begin{equation}\n\\int_0^1 x^2 dx\n\\end{equation}\n", encoding="utf-8")
    return path
