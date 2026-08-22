#!/usr/bin/env python3
"""Generate all static dummy files and test fixtures in structured subfolders.

Creates:
  - tests/dummies/canonical/ -> Standard fixtures (PDF, DOCX, PPTX, XLSX, ODT, multipage, etc.)
  - tests/dummies/formats/   -> All format conversion dummy files (dummy_file.<ext>) for every registered extension

Run this script once to pre-generate all required document formats so pytest
runs instantly without rebuilding temporary files during test runs.
"""

from __future__ import annotations

import bz2
import gzip
import io
import json
import os
import shutil
import sys
import tarfile
import zipfile
from email.message import EmailMessage
from pathlib import Path

# Add project root to sys.path so DEFAULT_REGISTRY can be imported
ROOT_DIR = Path(__file__).resolve().parent.parent.parent
sys.path.insert(0, str(ROOT_DIR))

from src.registry import DEFAULT_REGISTRY

DUMMIES_DIR = Path(__file__).resolve().parent
CANONICAL_DIR = DUMMIES_DIR / "canonical"
FORMATS_DIR = DUMMIES_DIR / "formats"

VERBATIM_MARKER = "VERBATIM_TEST_CONTENT_12345"
HEADER = "Field Notes - Chapter Overview"
FOOTER = "Confidential Draft"
PAGE_COUNT = 6


def generate_base_strings() -> None:
    """Ensure baseline strings file exists."""
    strings_path = DUMMIES_DIR / "strings.txt"
    strings_path.write_text(
        f"Book the First—Recalled to Life\n\nRecord #10492: Amount = $94,821.50\n\n$E = mc^2$\n\n{VERBATIM_MARKER}\n",
        encoding="utf-8",
    )
    print("  [✓] strings.txt")


def generate_canonical_fixtures() -> None:
    """Generate all fixtures used by tests/conftest.py and integration suites in tests/dummies/canonical/."""
    CANONICAL_DIR.mkdir(parents=True, exist_ok=True)

    # 1. Plain text & simple markup
    (CANONICAL_DIR / "readme.txt").write_text("Plain text content.\n", encoding="utf-8")
    (CANONICAL_DIR / "notes.md").write_text("# Notes\n\nSome markdown content here.\n", encoding="utf-8")
    (CANONICAL_DIR / "data.csv").write_text("name,age\nalice,30\nbob,25\n", encoding="utf-8")
    (CANONICAL_DIR / "config.json").write_text('{"host": "localhost", "port": 8080}', encoding="utf-8")
    (CANONICAL_DIR / "page.html").write_text("<html><body><article><h1>Title</h1><p>Body text here.</p></article></body></html>", encoding="utf-8")
    (CANONICAL_DIR / "sample.rtf").write_text(r"{\rtf1\ansi Hello {\b bold} RTF.}", encoding="utf-8")
    (CANONICAL_DIR / "sample.eml").write_text("From: a@example.com\nTo: b@example.com\nSubject: Hello\n\nBody text.\n", encoding="utf-8")
    (CANONICAL_DIR / "sample.srt").write_text("1\n00:00:01,000 --> 00:00:03,000\nHello subtitle\n", encoding="utf-8")
    (CANONICAL_DIR / "sample.tex").write_text("\\section{Introduction}\nSome body text here.\n", encoding="utf-8")
    (CANONICAL_DIR / "math.tex").write_text("\\section{Derivation}\nThe energy is $E = mc^2$.\n\\begin{equation}\n\\int_0^1 x^2 dx\n\\end{equation}\n", encoding="utf-8")

    # 2. PDF (Single-page & Multipage)
    try:
        import pymupdf

        doc = pymupdf.open()
        page = doc.new_page()
        page.insert_text((72, 72), "Hello from tokens.md")
        pdf_path = CANONICAL_DIR / "sample.pdf"
        doc.save(str(pdf_path))
        doc.close()

        # AZW4 wrapper
        shutil.copyfile(pdf_path, CANONICAL_DIR / "sample.azw4")

        # Multipage boilerplate PDF (6 pages)
        doc_multi = pymupdf.open()
        for i in range(1, PAGE_COUNT + 1):
            page_m = doc_multi.new_page()
            page_m.insert_text((50, 40), HEADER)
            page_m.insert_text((50, 100), f"Unique body content for page {i} about topic number {i}.")
            page_m.insert_text((50, 160), f"Additional distinct paragraph for page {i} with more words here.")
            page_m.insert_text((280, 780), str(i))
            page_m.insert_text((50, 800), FOOTER)
        multi_path = CANONICAL_DIR / "multipage.pdf"
        doc_multi.save(str(multi_path))
        doc_multi.close()
        shutil.copyfile(multi_path, CANONICAL_DIR / "report.pdf")
        print("  [✓] canonical/ (sample.pdf, sample.azw4, multipage.pdf, report.pdf)")
    except ImportError:
        print("  [!] Warning: pymupdf not found, skipping PDF generation.")

    # 3. DOCX (Standard, Headed, Math)
    try:
        import docx

        # Standard docx
        doc_std = docx.Document()
        doc_std.add_paragraph("Dear friend,")
        doc_std.add_paragraph("This is a letter.")
        doc_std.save(str(CANONICAL_DIR / "letter.docx"))
        shutil.copyfile(CANONICAL_DIR / "letter.docx", CANONICAL_DIR / "sample.docx")

        # Headed docx
        doc_h = docx.Document()
        doc_h.add_heading("Chapter One", level=1)
        doc_h.add_paragraph("Hello from python-docx.")
        doc_h.save(str(CANONICAL_DIR / "headed.docx"))

        # Math docx
        doc_m = docx.Document()
        doc_m.add_paragraph("Solve for x:")
        math_path = CANONICAL_DIR / "math.docx"
        doc_m.save(str(math_path))

        # Inject OMML: x^2 + 1
        omml = "<m:sSup><m:e><m:r><m:t>x</m:t></m:r></m:e><m:sup><m:r><m:t>2</m:t></m:r></m:sup></m:sSup><m:r><m:t>+1</m:t></m:r>"
        tmp_m = math_path.with_suffix(".tmp.docx")
        with zipfile.ZipFile(math_path) as zin, zipfile.ZipFile(tmp_m, "w") as zout:
            for item in zin.infolist():
                data = zin.read(item.filename)
                if item.filename == "word/document.xml":
                    xml = data.decode("utf-8")
                    omml_xml = f'<w:p><m:oMathPara xmlns:m="http://schemas.openxmlformats.org/officeDocument/2006/math">{omml}</m:oMathPara></w:p>'
                    xml = xml.replace("</w:p>", "</w:p>" + omml_xml, 1)
                    data = xml.encode("utf-8")
                zout.writestr(item, data)
        shutil.move(tmp_m, math_path)
        print("  [✓] canonical/ (letter.docx, sample.docx, headed.docx, math.docx)")
    except ImportError:
        print("  [!] Warning: python-docx not found, skipping DOCX generation.")

    # 4. PPTX
    try:
        import pptx

        prs = pptx.Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Hello slides"
        slide.placeholders[1].text = "Body text here."
        prs.save(str(CANONICAL_DIR / "deck.pptx"))
        shutil.copyfile(CANONICAL_DIR / "deck.pptx", CANONICAL_DIR / "sample.pptx")
        print("  [✓] canonical/ (deck.pptx, sample.pptx)")
    except ImportError:
        print("  [!] Warning: python-pptx not found, skipping PPTX generation.")

    # 5. XLSX
    try:
        import openpyxl

        wb = openpyxl.Workbook()
        ws = wb.active
        assert ws is not None
        ws.append(["name", "age"])
        ws.append(["alice", 30])
        wb.save(str(CANONICAL_DIR / "data.xlsx"))
        shutil.copyfile(CANONICAL_DIR / "data.xlsx", CANONICAL_DIR / "sample.xlsx")
        print("  [✓] canonical/ (data.xlsx, sample.xlsx)")
    except ImportError:
        print("  [!] Warning: openpyxl not found, skipping XLSX generation.")

    # 6. ODT
    try:
        from odf.opendocument import OpenDocumentText
        from odf.teletype import addTextToElement
        from odf.text import P

        doc_odt = OpenDocumentText()
        p = P()
        addTextToElement(p, "Hello from ODF")
        doc_odt.text.addElement(p)
        doc_odt.save(str(CANONICAL_DIR / "sample.odt"))
        print("  [✓] canonical/ (sample.odt)")
    except ImportError:
        print("  [!] Warning: odfpy not found, skipping ODT generation.")


def generate_format_dummies() -> None:
    """Generate dummy_file.<ext> for every registered format under tests/dummies/formats/."""
    FORMATS_DIR.mkdir(parents=True, exist_ok=True)
    corpus_text = (DUMMIES_DIR / "strings.txt").read_text(encoding="utf-8")

    extensions = sorted(DEFAULT_REGISTRY.extensions())

    for ext in extensions:
        dummy_path = FORMATS_DIR / f"dummy_file{ext}"

        if ext in {".txt", ".log", ".rst", ".markdown", ".md", ".mdx"}:
            dummy_path.write_text(corpus_text, encoding="utf-8")
        elif ext == ".csv":
            dummy_path.write_text(f"ID,Description,Amount,Status\n10492,Book the First—Recalled to Life,94821.50,{VERBATIM_MARKER}\n99012,Item Record 99012,17031.48,PROCESSED\n", encoding="utf-8")
        elif ext == ".tsv":
            dummy_path.write_text(f"ID\tDescription\tAmount\tStatus\n10492\tBook the First—Recalled to Life\t94821.50\t{VERBATIM_MARKER}\n99012\tItem Record 99012\t17031.48\tPROCESSED\n", encoding="utf-8")
        elif ext == ".json":
            dummy_path.write_text(json.dumps({"title": "Book the First—Recalled to Life", "marker": VERBATIM_MARKER, "equation": "E = mc^2", "records": [{"id": 10492, "amount": 94821.50, "token_count": 4291}, {"id": 99012, "quantity": 852, "price": 19.99}]}, indent=2), encoding="utf-8")
        elif ext in {".yaml", ".yml"}:
            dummy_path.write_text(f"chapter: Book the First—Recalled to Life\nmarker: {VERBATIM_MARKER}\nequation: E = mc^2\nrecords:\n  - id: 10492\n    amount: 94821.50\n    tokens: 4291\n", encoding="utf-8")
        elif ext == ".toml":
            dummy_path.write_text(f"title = 'Book the First—Recalled to Life'\nmarker = '{VERBATIM_MARKER}'\nequation = 'E = mc^2'\n[records]\nid = 10492\namount = 94821.50\n", encoding="utf-8")
        elif ext == ".ini":
            dummy_path.write_text(f"[Section]\nchapter = Book the First—Recalled to Life\nmarker = {VERBATIM_MARKER}\nrecord = 10492\n", encoding="utf-8")
        elif ext in {".html", ".htm"}:
            dummy_path.write_text(f"<html><body><h1>Book the First—Recalled to Life</h1><p>Record #10492: Amount = $94,821.50</p><p>{VERBATIM_MARKER}</p></body></html>", encoding="utf-8")
        elif ext == ".xml":
            dummy_path.write_text(f'<?xml version="1.0"?><root><chapter>Book the First—Recalled to Life</chapter><data>Record #10492</data><equation>E=mc^2</equation><marker>{VERBATIM_MARKER}</marker></root>', encoding="utf-8")
        elif ext == ".srt":
            dummy_path.write_text(f"1\n00:00:01,000 --> 00:00:04,000\nBook the First—Recalled to Life\n\n2\n00:00:04,500 --> 00:00:08,000\nRecord #10492: Amount = $94,821.50\n\n3\n00:00:08,500 --> 00:00:12,000\n{VERBATIM_MARKER}\n\n", encoding="utf-8")
        elif ext == ".vtt":
            dummy_path.write_text(f"WEBVTT\n\n1\n00:00:01.000 --> 00:00:04.000\nBook the First—Recalled to Life\n\n2\n00:00:04.500 --> 00:00:08.000\nRecord #10492: Amount = $94,821.50\n\n3\n00:00:08.500 --> 00:00:12.000\n{VERBATIM_MARKER}\n\n", encoding="utf-8")
        elif ext == ".tex":
            dummy_path.write_text(f"\\documentclass{{article}}\n\\begin{{document}}\n\\section*{{Book the First---Recalled to Life}}\nRecord #10492: Amount = \\$94,821.50\n$E = mc^2$\n{VERBATIM_MARKER}\n\\end{{document}}\n", encoding="utf-8")
        elif ext == ".ipynb":
            notebook = {"cells": [{"cell_type": "markdown", "source": ["# Book the First—Recalled to Life\n", "E = mc^2"]}, {"cell_type": "code", "source": ["# Record #10492\n", f"marker = '{VERBATIM_MARKER}'"]}], "metadata": {}, "nbformat": 4, "nbformat_minor": 2}
            dummy_path.write_text(json.dumps(notebook), encoding="utf-8")
        elif ext == ".eml":
            msg = EmailMessage()
            msg["From"] = "sender@example.com"
            msg["To"] = "recv@example.com"
            msg["Subject"] = "Book the First—Recalled to Life"
            msg.set_content(f"Record #10492: Amount = $94,821.50\nE = mc^2\n{VERBATIM_MARKER}")
            dummy_path.write_bytes(msg.as_bytes())
        elif ext == ".rtf":
            dummy_path.write_text(f"{{\\rtf1\\ansi Book the First---Recalled to Life\\par Record #10492\\par {VERBATIM_MARKER}}}", encoding="utf-8")
        elif ext == ".pdf":
            import pymupdf
            doc = pymupdf.open()
            page = doc.new_page()
            page.insert_text((50, 50), "Book the First - Recalled to Life")
            page.insert_text((50, 80), "Record #10492: Amount = $94,821.50")
            page.insert_text((50, 110), VERBATIM_MARKER)
            doc.save(str(dummy_path))
            doc.close()
        elif ext == ".docx":
            import docx
            doc = docx.Document()
            doc.add_heading("Book the First—Recalled to Life", level=1)
            doc.add_paragraph("Record #10492: Amount = $94,821.50 | Token Count = 4291")
            doc.add_paragraph(VERBATIM_MARKER)
            doc.save(str(dummy_path))
        elif ext == ".xlsx":
            import openpyxl
            wb = openpyxl.Workbook()
            ws = wb.active
            assert ws is not None
            ws.append(["Chapter", "Record Info", "Marker"])
            ws.append(["Book the First—Recalled to Life", "Record #10492: Amount = $94,821.50", VERBATIM_MARKER])
            wb.save(str(dummy_path))
        elif ext == ".pptx":
            import pptx
            prs = pptx.Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            txBox = slide.shapes.add_textbox(0, 0, 400, 300)
            txBox.text_frame.text = f"Book the First—Recalled to Life\nRecord #10492\n{VERBATIM_MARKER}"
            prs.save(str(dummy_path))
        elif ext == ".odt":
            from odf import text as odftext
            from odf.opendocument import OpenDocumentText
            doc = OpenDocumentText()
            doc.text.addElement(odftext.P(text="Book the First—Recalled to Life"))
            doc.text.addElement(odftext.P(text="Record #10492: Amount = $94,821.50"))
            doc.text.addElement(odftext.P(text=VERBATIM_MARKER))
            doc.save(str(dummy_path))
        elif ext == ".ods":
            from odf import text as odftext
            from odf.opendocument import OpenDocumentSpreadsheet
            from odf.table import Table as OdfTable, TableCell, TableRow
            doc = OpenDocumentSpreadsheet()
            tbl = OdfTable()
            tr = TableRow()
            tc = TableCell()
            tc.addElement(odftext.P(text=f"Book the First—Recalled to Life {VERBATIM_MARKER}"))
            tr.addElement(tc)
            tbl.addElement(tr)
            doc.spreadsheet.addElement(tbl)
            doc.save(str(dummy_path))
        elif ext == ".odp":
            from odf import text as odftext
            from odf.draw import Frame, Page, TextBox
            from odf.opendocument import OpenDocumentPresentation
            doc = OpenDocumentPresentation()
            page = Page(masterpagename="Standard")
            frame = Frame()
            tb = TextBox()
            tb.addElement(odftext.P(text=f"Book the First—Recalled to Life {VERBATIM_MARKER}"))
            frame.addElement(tb)
            page.addElement(frame)
            doc.presentation.addElement(page)
            doc.save(str(dummy_path))
        elif ext == ".fb2":
            dummy_path.write_text(f'<?xml version="1.0" encoding="utf-8"?><FictionBook xmlns="http://www.gribuser.ru/xml/fictionbook/2.0"><body><section><title><p>Book the First—Recalled to Life</p></title><p>{VERBATIM_MARKER}</p></section></body></FictionBook>', encoding="utf-8")
        elif ext == ".epub":
            with zipfile.ZipFile(dummy_path, "w") as zf:
                zf.writestr("mimetype", "application/epub+zip")
                zf.writestr("META-INF/container.xml", '<?xml version="1.0"?><container version="1.0" xmlns="urn:oasis:names:tc:opendocument:xmlns:container"><rootfiles><rootfile full-path="content.opf" media-type="application/oebps-package+xml"/></rootfiles></container>')
                zf.writestr("content.opf", '<?xml version="1.0"?><package xmlns="http://www.idpf.org/2007/opf" unique-identifier="BookID" version="2.0"><metadata></metadata><manifest><item id="chap1" href="chap1.xhtml" media-type="application/xhtml+xml"/></manifest><spine><itemref idref="chap1"/></spine></package>')
                zf.writestr("chap1.xhtml", f'<?xml version="1.0"?><html xmlns="http://www.w3.org/1999/xhtml"><body><h1>Book the First—Recalled to Life</h1><p>{VERBATIM_MARKER}</p></body></html>')
        elif ext in {".xps", ".oxps"}:
            with zipfile.ZipFile(dummy_path, "w") as zf:
                zf.writestr("[Content_Types].xml", '<?xml version="1.0"?><Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types"><Default Extension="fdseq" ContentType="application/vnd.ms-package.xps-fixeddocumentsequence+xml"/><Default Extension="fpage" ContentType="application/vnd.ms-package.xps-fixedpage+xml"/></Types>')
                zf.writestr("FixedDocumentSequence.fdseq", '<?xml version="1.0"?><FixedDocumentSequence xmlns="http://schemas.microsoft.com/xps/2005/06"><DocumentReference Source="Documents/1/FixedDocument.fdoc"/></FixedDocumentSequence>')
                zf.writestr("Documents/1/FixedDocument.fdoc", '<?xml version="1.0"?><FixedDocument xmlns="http://schemas.microsoft.com/xps/2005/06"><PageContent Source="1.fpage"/></FixedDocument>')
                zf.writestr("Documents/1/Pages/1.fpage", f'<?xml version="1.0"?><FixedPage xmlns="http://schemas.microsoft.com/xps/2005/06" Width="816" Height="1056"><Glyphs UnicodeString="Book the First {VERBATIM_MARKER}"/></FixedPage>')
        elif ext == ".cbz":
            with zipfile.ZipFile(dummy_path, "w") as zf:
                zf.writestr("page1.txt", f"Book the First—Recalled to Life\n{VERBATIM_MARKER}\n")
        elif ext == ".zip":
            with zipfile.ZipFile(dummy_path, "w") as zf:
                zf.writestr("chapter1.txt", f"Book the First—Recalled to Life\n{VERBATIM_MARKER}\n")
        elif ext in {".tar", ".tgz"}:
            with tarfile.open(dummy_path, "w:gz" if ext == ".tgz" else "w") as tf:
                data = f"Book the First—Recalled to Life\n{VERBATIM_MARKER}\n".encode()
                ti = tarfile.TarInfo("chapter1.txt")
                ti.size = len(data)
                tf.addfile(ti, io.BytesIO(data))
        elif ext == ".gz":
            with gzip.open(dummy_path, "wb") as gf:
                gf.write(f"Book the First—Recalled to Life\n{VERBATIM_MARKER}\n".encode())
        elif ext == ".bz2":
            with bz2.open(dummy_path, "wb") as bf:
                bf.write(f"Book the First—Recalled to Life\n{VERBATIM_MARKER}\n".encode())
        elif ext in {".msg", ".azw3", ".azw4", ".mobi"}:
            dummy_path.write_bytes(b"Dummy binary data")
        else:
            dummy_path.write_text(corpus_text, encoding="utf-8")

    print(f"  [✓] formats/ (Generated {len(extensions)} format dummy files)")


def main() -> None:
    print(f"Generating structured dummy test files in: {DUMMIES_DIR}")
    generate_base_strings()
    generate_canonical_fixtures()
    generate_format_dummies()
    print("\n✓ All dummy test files successfully generated into subfolders.")


if __name__ == "__main__":
    main()
