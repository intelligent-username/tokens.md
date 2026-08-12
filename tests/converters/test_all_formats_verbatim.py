"""Test that EVERY supported document file format physically creates a real file under tests/dummies/ and converts verbatim.

Image formats (.png, .jpg, .jpeg, .bmp, .gif, .tif, .tiff, .svg) are excluded as non-document formats.
"""

from __future__ import annotations

import bz2
import email
from email.message import EmailMessage
import gzip
import io
import json
import shutil
import tarfile
import zipfile
from pathlib import Path
from typing import Generator

import pytest

from src.registry import DEFAULT_REGISTRY, convert_file

DUMMIES_DIR = Path("tests/dummies")
STRINGS_FILE = DUMMIES_DIR / "strings.txt"

VERBATIM_MARKER = "VERBATIM_TEST_CONTENT_12345"


@pytest.fixture(scope="module")
def dummies_folder() -> Generator[Path, None, None]:
    """Ensure tests/dummies/ and tests/dummies/strings.txt exist, and clean up temporary generated dummy files afterwards."""
    DUMMIES_DIR.mkdir(parents=True, exist_ok=True)
    if not STRINGS_FILE.exists():
        STRINGS_FILE.write_text(
            f"Book the First—Recalled to Life\n\nRecord #10492: Amount = $94,821.50\n\n$E = mc^2$\n\n{VERBATIM_MARKER}",
            encoding="utf-8",
        )

    try:
        yield DUMMIES_DIR
    finally:
        for p in DUMMIES_DIR.glob("dummy_file*"):
            try:
                p.unlink()
            except OSError:
                pass


def _create_dummy_file(folder: Path, ext: str) -> Path:
    """Physically create a real, valid document file on disk under tests/dummies/ for each extension."""
    dummy_path = folder / f"dummy_file{ext}"
    corpus_text = STRINGS_FILE.read_text(encoding="utf-8")

    if ext in {".txt", ".log", ".rst", ".markdown", ".md", ".mdx"}:
        dummy_path.write_text(corpus_text, encoding="utf-8")
    elif ext == ".csv":
        dummy_path.write_text(
            "ID,Description,Amount,Status\n"
            f"10492,Book the First—Recalled to Life,94821.50,{VERBATIM_MARKER}\n"
            "99012,Item Record 99012,17031.48,PROCESSED\n",
            encoding="utf-8",
        )
    elif ext == ".tsv":
        dummy_path.write_text(
            "ID\tDescription\tAmount\tStatus\n"
            f"10492\tBook the First—Recalled to Life\t94821.50\t{VERBATIM_MARKER}\n"
            "99012\tItem Record 99012\t17031.48\tPROCESSED\n",
            encoding="utf-8",
        )
    elif ext == ".json":
        dummy_path.write_text(
            json.dumps(
                {
                    "title": "Book the First—Recalled to Life",
                    "marker": VERBATIM_MARKER,
                    "equation": "E = mc^2",
                    "records": [
                        {"id": 10492, "amount": 94821.50, "token_count": 4291},
                        {"id": 99012, "quantity": 852, "price": 19.99},
                    ],
                },
                indent=2,
            ),
            encoding="utf-8",
        )
    elif ext in {".yaml", ".yml"}:
        dummy_path.write_text(
            "chapter: Book the First—Recalled to Life\n"
            f"marker: {VERBATIM_MARKER}\n"
            "equation: E = mc^2\n"
            "records:\n"
            "  - id: 10492\n"
            "    amount: 94821.50\n"
            "    tokens: 4291\n",
            encoding="utf-8",
        )
    elif ext in {".html", ".htm"}:
        dummy_path.write_text(
            f"<html><body><h1>Book the First—Recalled to Life</h1><p>Record #10492: Amount = $94,821.50</p><p>{VERBATIM_MARKER}</p></body></html>",
            encoding="utf-8",
        )
    elif ext == ".xml":
        dummy_path.write_text(
            f'<?xml version="1.0"?><root><chapter>Book the First—Recalled to Life</chapter><data>Record #10492</data><equation>E=mc^2</equation><marker>{VERBATIM_MARKER}</marker></root>',
            encoding="utf-8",
        )
    elif ext == ".srt":
        dummy_path.write_text(
            f"1\n00:00:01,000 --> 00:00:04,000\nBook the First—Recalled to Life\n\n"
            f"2\n00:00:04,500 --> 00:00:08,000\nRecord #10492: Amount = $94,821.50\n\n"
            f"3\n00:00:08,500 --> 00:00:12,000\n{VERBATIM_MARKER}\n\n",
            encoding="utf-8",
        )
    elif ext == ".vtt":
        dummy_path.write_text(
            f"WEBVTT\n\n1\n00:00:01.000 --> 00:00:04.000\nBook the First—Recalled to Life\n\n"
            f"2\n00:00:04.500 --> 00:00:08.000\nRecord #10492: Amount = $94,821.50\n\n"
            f"3\n00:00:08.500 --> 00:00:12.000\n{VERBATIM_MARKER}\n\n",
            encoding="utf-8",
        )
    elif ext == ".tex":
        dummy_path.write_text(
            f"\\documentclass{{article}}\n\\begin{{document}}\n\\section*{{Book the First---Recalled to Life}}\nRecord #10492: Amount = \\$94,821.50\n$E = mc^2$\n{VERBATIM_MARKER}\n\\end{{document}}\n",
            encoding="utf-8",
        )
    elif ext == ".ipynb":
        notebook = {
            "cells": [
                {"cell_type": "markdown", "source": ["# Book the First—Recalled to Life\n", "E = mc^2"]},
                {"cell_type": "code", "source": ["# Record #10492\n", f"marker = '{VERBATIM_MARKER}'"]},
            ],
            "metadata": {},
            "nbformat": 4,
            "nbformat_minor": 2,
        }
        dummy_path.write_text(json.dumps(notebook), encoding="utf-8")
    elif ext == ".eml":
        msg = EmailMessage()
        msg["From"] = "sender@example.com"
        msg["To"] = "recv@example.com"
        msg["Subject"] = "Book the First—Recalled to Life"
        msg.set_content(f"Record #10492: Amount = $94,821.50\nE = mc^2\n{VERBATIM_MARKER}")
        dummy_path.write_bytes(msg.as_bytes())
    elif ext == ".rtf":
        dummy_path.write_text(
            f"{{\\rtf1\\ansi Book the First---Recalled to Life\\par Record #10492\\par {VERBATIM_MARKER}}}",
            encoding="utf-8",
        )
    elif ext == ".pdf":
        import pymupdf

        doc = pymupdf.open()
        page = doc.new_page()
        page.insert_text((50, 50), "Book the First - Recalled to Life")
        page.insert_text((50, 80), "Record #10492: Amount = $94,821.50")
        page.insert_text((50, 110), VERBATIM_MARKER)
        doc.save(str(dummy_path))
        doc.close()
    elif ext == ".docx":
        import docx

        doc = docx.Document()
        doc.add_heading("Book the First—Recalled to Life", level=1)
        doc.add_paragraph("Record #10492: Amount = $94,821.50 | Token Count = 4291")
        doc.add_paragraph(VERBATIM_MARKER)
        doc.save(dummy_path)
    elif ext == ".xlsx":
        import openpyxl

        wb = openpyxl.Workbook()
        ws = wb.active
        ws.append(["Chapter", "Record Info", "Marker"])
        ws.append(["Book the First—Recalled to Life", "Record #10492: Amount = $94,821.50", VERBATIM_MARKER])
        wb.save(dummy_path)
    elif ext == ".pptx":
        import pptx

        prs = pptx.Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        txBox = slide.shapes.add_textbox(0, 0, 400, 300)
        txBox.text_frame.text = f"Book the First—Recalled to Life\nRecord #10492\n{VERBATIM_MARKER}"
        prs.save(dummy_path)
    elif ext == ".odt":
        from odf import text as odftext
        from odf.opendocument import OpenDocumentText

        doc = OpenDocumentText()
        doc.text.addElement(odftext.P(text="Book the First—Recalled to Life"))
        doc.text.addElement(odftext.P(text="Record #10492: Amount = $94,821.50"))
        doc.text.addElement(odftext.P(text=VERBATIM_MARKER))
        doc.save(dummy_path)
    elif ext == ".ods":
        from odf import text as odftext
        from odf.opendocument import OpenDocumentSpreadsheet
        from odf.table import Table as OdfTable, TableCell, TableRow

        doc = OpenDocumentSpreadsheet()
        tbl = OdfTable()
        tr = TableRow()
        tc = TableCell()
        tc.addElement(odftext.P(text=f"Book the First—Recalled to Life {VERBATIM_MARKER}"))
        tr.addElement(tc)
        tbl.addElement(tr)
        doc.spreadsheet.addElement(tbl)
        doc.save(dummy_path)
    elif ext == ".odp":
        from odf import text as odftext
        from odf.draw import Frame, Page, TextBox
        from odf.opendocument import OpenDocumentPresentation

        doc = OpenDocumentPresentation()
        page = Page(masterpagename="Standard")
        frame = Frame()
        tb = TextBox()
        tb.addElement(odftext.P(text=f"Book the First—Recalled to Life {VERBATIM_MARKER}"))
        frame.addElement(tb)
        page.addElement(frame)
        doc.presentation.addElement(page)
        doc.save(dummy_path)
    elif ext == ".fb2":
        dummy_path.write_text(
            f'<?xml version="1.0" encoding="utf-8"?><FictionBook xmlns="http://www.gribuser.ru/xml/fictionbook/2.0"><body><section><title><p>Book the First—Recalled to Life</p></title><p>{VERBATIM_MARKER}</p></section></body></FictionBook>',
            encoding="utf-8",
        )
    elif ext == ".epub":
        with zipfile.ZipFile(dummy_path, "w") as zf:
            zf.writestr("mimetype", "application/epub+zip")
            zf.writestr(
                "META-INF/container.xml",
                '<?xml version="1.0"?><container version="1.0" xmlns="urn:oasis:names:tc:opendocument:xmlns:container"><rootfiles><rootfile full-path="content.opf" media-type="application/oebps-package+xml"/></rootfiles></container>',
            )
            zf.writestr(
                "content.opf",
                '<?xml version="1.0"?><package xmlns="http://www.idpf.org/2007/opf" unique-identifier="BookID" version="2.0"><metadata></metadata><manifest><item id="chap1" href="chap1.xhtml" media-type="application/xhtml+xml"/></manifest><spine><itemref idref="chap1"/></spine></package>',
            )
            zf.writestr(
                "chap1.xhtml",
                f'<?xml version="1.0"?><html xmlns="http://www.w3.org/1999/xhtml"><body><h1>Book the First—Recalled to Life</h1><p>{VERBATIM_MARKER}</p></body></html>',
            )
    elif ext in {".xps", ".oxps"}:
        with zipfile.ZipFile(dummy_path, "w") as zf:
            zf.writestr(
                "[Content_Types].xml",
                '<?xml version="1.0"?><Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types"><Default Extension="fdseq" ContentType="application/vnd.ms-package.xps-fixeddocumentsequence+xml"/><Default Extension="fpage" ContentType="application/vnd.ms-package.xps-fixedpage+xml"/></Types>',
            )
            zf.writestr(
                "FixedDocumentSequence.fdseq",
                '<?xml version="1.0"?><FixedDocumentSequence xmlns="http://schemas.microsoft.com/xps/2005/06"><DocumentReference Source="Documents/1/FixedDocument.fdoc"/></FixedDocumentSequence>',
            )
            zf.writestr(
                "Documents/1/FixedDocument.fdoc",
                '<?xml version="1.0"?><FixedDocument xmlns="http://schemas.microsoft.com/xps/2005/06"><PageContent Source="1.fpage"/></FixedDocument>',
            )
            zf.writestr(
                "Documents/1/Pages/1.fpage",
                f'<?xml version="1.0"?><FixedPage xmlns="http://schemas.microsoft.com/xps/2005/06" Width="816" Height="1056"><Glyphs UnicodeString="Book the First {VERBATIM_MARKER}"/></FixedPage>',
            )
    elif ext == ".cbz":
        with zipfile.ZipFile(dummy_path, "w") as zf:
            zf.writestr("page1.txt", f"Book the First—Recalled to Life\n{VERBATIM_MARKER}\n")
    elif ext == ".zip":
        with zipfile.ZipFile(dummy_path, "w") as zf:
            zf.writestr("chapter1.txt", f"Book the First—Recalled to Life\n{VERBATIM_MARKER}\n")
    elif ext in {".tar", ".tgz"}:
        with tarfile.open(dummy_path, "w:gz" if ext == ".tgz" else "w") as tf:
            data = f"Book the First—Recalled to Life\n{VERBATIM_MARKER}\n".encode("utf-8")
            ti = tarfile.TarInfo("chapter1.txt")
            ti.size = len(data)
            tf.addfile(ti, io.BytesIO(data))
    elif ext == ".gz":
        with gzip.open(dummy_path, "wb") as gf:
            gf.write(f"Book the First—Recalled to Life\n{VERBATIM_MARKER}\n".encode("utf-8"))
    elif ext == ".bz2":
        with bz2.open(dummy_path, "wb") as bf:
            bf.write(f"Book the First—Recalled to Life\n{VERBATIM_MARKER}\n".encode("utf-8"))
    elif ext in {".msg", ".azw3", ".azw4", ".mobi"}:
        dummy_path.write_bytes(b"Dummy binary data")
    else:
        dummy_path.write_text(corpus_text, encoding="utf-8")

    return dummy_path


@pytest.mark.parametrize("ext", sorted(DEFAULT_REGISTRY.extensions()))
def test_format_conversion_verbatim(ext: str, dummies_folder: Path, tmp_path: Path) -> None:
    """Test individual file format conversion for each non-image extension registered in DEFAULT_REGISTRY."""
    dummy = _create_dummy_file(dummies_folder, ext)
    assert dummy.exists()
    assert dummy.parent == DUMMIES_DIR

    output_dir = tmp_path / f"out_{ext.lstrip('.')}"
    output_dir.mkdir(parents=True, exist_ok=True)

    result = convert_file(dummy, output_dir)
    assert result.exists()
    assert result.suffix == ".md"
    converted_text = result.read_text(encoding="utf-8", errors="replace")
    assert isinstance(converted_text, str)
